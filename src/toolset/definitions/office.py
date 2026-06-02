import os
from typing import List, Dict, Any, Type, Optional
from pydantic import BaseModel, Field
from toolset.base import BaseTool


# --- 1. EXCEL 工具 ---
class ExcelWriteArgs(BaseModel):
    path: str = Field(description="保存 Excel 的绝对路径。禁止在文件名中添加随机数或时间戳，确保任务周期内路径唯一。")
    data: List[Dict[str, Any]] = Field(description="列表形式的 JSON 数据，例如：[{'公司': 'OpenAI', '估值': '800亿'}]")
    sheet_name: str = Field(description="工作表名称", default="Sheet1")


class ExcelWriteTool(BaseTool):
    """机械化 Excel 写入工具"""

    name: str = "excel_write"
    kit: str = "Office"
    fc_hidden: bool = True  # [Round 10] Use excel_op(action="write") instead
    description: str = "Save structured JSON data as an Excel (.xlsx) file. Input must be a JSON array of objects with consistent keys."
    domain: str = "craft"
    args_schema: Type[BaseModel] = ExcelWriteArgs

    async def run(self, **kwargs) -> str:
        try:
            import pandas as pd
        except ImportError:
            return "Error: 'pandas' or 'openpyxl' not installed. Please run 'pip install pandas openpyxl'."

        path = kwargs.get("path")
        data = kwargs.get("data", [])
        sheet_name = kwargs.get("sheet_name", "Sheet1")

        try:
            df = pd.DataFrame(data)
            # Ensure parent directory exists (orchestrator handles path validity, but tool layer still creates physically)
            # 确保父目录存在 (编排器会负责路径合法性，但工具层依然做物理创建)
            os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
            df.to_excel(path, index=False, sheet_name=sheet_name)
            return f"Successfully saved {len(data)} rows to Excel. [RESULT_PATH: {os.path.abspath(path)}]"
        except Exception as e:
            return f"Excel Write Error: {str(e)}"


class ExcelReadArgs(BaseModel):
    path: str = Field(description="Excel 文件路径")


class ExcelReadTool(BaseTool):
    """机械化 Excel 读取工具"""

    name: str = "excel_read"
    kit: str = "Office"
    fc_hidden: bool = True  # [Round 10] Use excel_op(action="read") instead
    description: str = (
        "Read an Excel file and return its content as CSV-style text. Use this to inspect spreadsheet data."
    )
    domain: str = "craft"
    args_schema: Type[BaseModel] = ExcelReadArgs

    async def run(self, **kwargs) -> str:
        try:
            import pandas as pd
        except ImportError:
            return "Error: 'pandas' or 'openpyxl' not installed. Please run 'pip install pandas openpyxl'."

        path = kwargs.get("path")
        try:
            df = pd.read_excel(path)
            # Return first 20 rows to prevent context overflow (orchestrator applies secondary truncation)
            # 返回前 20 条，防止上下文溢出（编排器会有二次截断，此处做初步控制）
            csv_preview = df.to_csv(index=False)
            return f"Excel Content Preview (Total {len(df)} rows):\n{csv_preview}"
        except Exception as e:
            return f"Excel Read Error: {str(e)}"


# --- 2. WORD (DOCX) 工具 ---
class DocxWriteArgs(BaseModel):
    path: str = Field(description="保存 Word (.docx) 的路径")
    markdown_content: str = Field(description="包含 Markdown 语法的内容，支持标题、列表、表格、图片")
    style: str = Field(
        default="default",
        description='文档样式。预设: "公文"(中国公文格式), "学术"(论文格式), "商务"(商务报告), "default"(通用专业格式)。'
        '也可传 JSON 字符串自定义，如: {"title_font":"黑体","title_size":22,"body_font":"仿宋","body_size":16,"line_spacing":30}',
    )


class DocxWriteTool(BaseTool):
    """通用 Word 文档生成工具 — 支持多种样式预设和自定义格式"""

    name: str = "office_docx_write"
    kit: str = "Office"
    description: str = (
        "将 Markdown 内容转换为格式化的 Word (.docx) 文档。"
        "【追加写入特性】：如果目标文件路径已存在，此工具将自动把新内容**追加**在文档末尾！这允许你分多次、分章节调用此工具来生成超长报告，而无需一次性输出全部内容。"
        "支持样式预设：公文（中国公文排版标准）、学术（论文格式）、商务（商务报告）、default（通用专业）。"
        "用户可在对话中指定字体、字号、行距等格式要求，通过 style 参数传入。"
    )
    domain: str = "craft"
    args_schema: Type[BaseModel] = DocxWriteArgs

    # === 样式预设 ===
    STYLE_PRESETS = {
        "公文": {
            "title_font": "华文中宋",
            "title_size": 22,
            "title_bold": True,
            "body_font": "仿宋",
            "body_size": 16,
            "body_bold": False,
            "level1_font": "黑体",
            "level2_font": "楷体",
            "level3_font": "仿宋",
            "level4_font": "仿宋",
            "line_spacing": 30,
            "first_line_indent_chars": 2,
            "title_align": "center",
            "margins": {"top": 3.7, "bottom": 3.5, "left": 2.8, "right": 2.6},
        },
        "学术": {
            "title_font": "黑体",
            "title_size": 18,
            "title_bold": True,
            "body_font": "宋体",
            "body_size": 12,
            "body_bold": False,
            "level1_font": "黑体",
            "level2_font": "黑体",
            "level3_font": "宋体",
            "level4_font": "宋体",
            "line_spacing": 22,
            "first_line_indent_chars": 2,
            "title_align": "center",
            "margins": {"top": 2.54, "bottom": 2.54, "left": 3.18, "right": 3.18},
        },
        "商务": {
            "title_font": "微软雅黑",
            "title_size": 18,
            "title_bold": True,
            "body_font": "微软雅黑",
            "body_size": 11,
            "body_bold": False,
            "level1_font": "微软雅黑",
            "level2_font": "微软雅黑",
            "level3_font": "微软雅黑",
            "level4_font": "微软雅黑",
            "line_spacing": 26,
            "first_line_indent_chars": 2,
            "title_align": "center",
            "margins": {"top": 2.54, "bottom": 2.54, "left": 2.54, "right": 2.54},
        },
        "default": {
            "title_font": "宋体",
            "title_size": 18,
            "title_bold": True,
            "body_font": "宋体",
            "body_size": 12,
            "body_bold": False,
            "level1_font": "黑体",
            "level2_font": "楷体",
            "level3_font": "宋体",
            "level4_font": "宋体",
            "line_spacing": 26,
            "first_line_indent_chars": 2,
            "title_align": "center",
            "margins": {"top": 2.54, "bottom": 2.54, "left": 3.18, "right": 3.18},
        },
    }

    async def run(self, **kwargs) -> str:
        try:
            from docx import Document
            from docx.shared import Pt, Inches, Cm
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.ns import qn
        except ImportError:
            return "Error: 'python-docx' not installed. Please run 'pip install python-docx'."

        path = kwargs.get("path")
        md = kwargs.get("markdown_content", "")
        style_str = kwargs.get("style", "default")

        try:
            import re

            # --- 解析样式 ---
            style = self._resolve_style(style_str)

            # --- 从 style 提取参数 ---
            title_font = style.get("title_font", "宋体")
            title_size = Pt(style.get("title_size", 18))
            title_bold = style.get("title_bold", True)
            body_font = style.get("body_font", "宋体")
            body_size = Pt(style.get("body_size", 12))
            _body_bold = style.get("body_bold", False)  # noqa: F841 — TODO: pass to set_font() calls
            l1_font = style.get("level1_font", "黑体")
            l2_font = style.get("level2_font", "楷体")
            l3_font = style.get("level3_font", body_font)
            l4_font = style.get("level4_font", body_font)
            line_sp = Pt(style.get("line_spacing", 26))
            indent_chars = style.get("first_line_indent_chars", 2)
            indent = Pt(body_size.pt * indent_chars) if indent_chars else None
            title_align_str = style.get("title_align", "center")
            title_align = WD_ALIGN_PARAGRAPH.CENTER if title_align_str == "center" else WD_ALIGN_PARAGRAPH.LEFT
            margins = style.get("margins", {"top": 2.54, "bottom": 2.54, "left": 3.18, "right": 3.18})

            is_append = False
            if os.path.exists(path):
                try:
                    doc = Document(path)
                    is_append = True
                except Exception:
                    doc = Document()
            else:
                doc = Document()

            def set_font(run, font_name, size, bold=False):
                run.font.size = size
                run.font.bold = bold
                run.font.name = font_name
                r = run._element
                rPr = r.find(qn("w:rPr"))
                if rPr is None:
                    rPr = r.makeelement(qn("w:rPr"), {})
                    r.insert(0, rPr)
                rFonts = rPr.find(qn("w:rFonts"))
                if rFonts is None:
                    rFonts = rPr.makeelement(qn("w:rFonts"), {})
                    rPr.insert(0, rFonts)
                rFonts.set(qn("w:eastAsia"), font_name)

            def set_para(
                p,
                alignment=WD_ALIGN_PARAGRAPH.JUSTIFY,
                line_spacing=line_sp,
                first_indent=None,
                space_before=Pt(0),
                space_after=Pt(0),
            ):
                pf = p.paragraph_format
                pf.alignment = alignment
                pf.line_spacing = line_spacing
                pf.space_before = space_before
                pf.space_after = space_after
                if first_indent is not None:
                    pf.first_line_indent = first_indent

            def add_paragraph(
                doc,
                text,
                font,
                size,
                bold=False,
                alignment=WD_ALIGN_PARAGRAPH.JUSTIFY,
                first_indent=None,
                space_before=Pt(0),
                space_after=Pt(0),
            ):
                p = doc.add_paragraph()
                set_para(
                    p,
                    alignment=alignment,
                    first_indent=first_indent,
                    space_before=space_before,
                    space_after=space_after,
                )
                run = p.add_run(text)
                set_font(run, font, size, bold)
                return p

            def add_rich_paragraph(doc, text, font, size, first_indent=None, bold_font=None):
                """处理含 **粗体** 的段落"""
                p = doc.add_paragraph()
                set_para(p, first_indent=first_indent)
                bold_pattern = re.compile(r"\*\*(.*?)\*\*")
                if bold_pattern.search(text):
                    last_end = 0
                    for m in bold_pattern.finditer(text):
                        if m.start() > last_end:
                            run = p.add_run(text[last_end : m.start()])
                            set_font(run, font, size)
                        run = p.add_run(m.group(1))
                        set_font(run, bold_font or font, size, bold=True)
                        last_end = m.end()
                    if last_end < len(text):
                        run = p.add_run(text[last_end:])
                        set_font(run, font, size)
                else:
                    run = p.add_run(text)
                    set_font(run, font, size)
                return p

            # === Parse Markdown ===
            # === 解析 Markdown ===
            lines = md.split("\n")
            i = 0
            is_first_heading = True

            while i < len(lines):
                line = lines[i].strip()
                if not line:
                    i += 1
                    continue

                # 1. Table
                # 1. 表格
                if line.startswith("|") and i + 1 < len(lines) and "|---" in lines[i + 1]:
                    headers = [c.strip() for c in line.split("|") if c.strip()]
                    table = doc.add_table(rows=1, cols=len(headers))
                    for idx, h in enumerate(headers):
                        cell = table.rows[0].cells[idx]
                        cell.text = ""
                        run = cell.paragraphs[0].add_run(h)
                        set_font(run, body_font, body_size, bold=True)
                        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                    i += 2
                    while i < len(lines) and lines[i].strip().startswith("|"):
                        row_data = [c.strip() for c in lines[i].split("|") if c.strip()]
                        if len(row_data) == len(headers):
                            row_cells = table.add_row().cells
                            for idx, val in enumerate(row_data):
                                row_cells[idx].text = ""
                                run = row_cells[idx].paragraphs[0].add_run(val)
                                set_font(run, body_font, body_size)
                        i += 1
                    doc.add_paragraph()
                    continue

                # 2. Image
                # 2. 图片
                image_match = re.match(r"!\[(.*?)\]\((.*?)\)", line)
                if image_match:
                    try:
                        doc.add_picture(image_match.group(2), width=Inches(4))
                        if image_match.group(1):
                            add_paragraph(
                                doc, image_match.group(1), body_font, Pt(10), alignment=WD_ALIGN_PARAGRAPH.CENTER
                            )
                    except Exception as img_e:
                        add_paragraph(doc, f"[图片插入失败: {img_e}]", body_font, body_size)
                    i += 1
                    continue

                # 3. Markdown heading
                # 3. Markdown 标题
                if line.startswith("# "):
                    text = line[2:].strip()
                    if is_first_heading and not is_append:
                        add_paragraph(
                            doc,
                            text,
                            title_font,
                            title_size,
                            bold=title_bold,
                            alignment=title_align,
                            space_before=Pt(12),
                            space_after=Pt(12),
                        )
                        is_first_heading = False
                    else:
                        add_paragraph(doc, text, l1_font, body_size, space_before=Pt(6), space_after=Pt(3))
                        is_first_heading = False
                elif line.startswith("## "):
                    add_paragraph(doc, line[3:].strip(), l1_font, body_size, space_before=Pt(6), space_after=Pt(3))
                elif line.startswith("### "):
                    add_paragraph(
                        doc,
                        line[4:].strip(),
                        l2_font,
                        body_size,
                        first_indent=indent,
                        space_before=Pt(3),
                        space_after=Pt(3),
                    )
                elif line.startswith("#### "):
                    add_paragraph(doc, line[5:].strip(), l3_font, body_size, first_indent=indent)
                # 4. Official document native hierarchy
                # 4. 公文原生层次
                elif re.match(r"^[一二三四五六七八九十]+、", line):
                    add_paragraph(doc, line, l1_font, body_size, space_before=Pt(6), space_after=Pt(3))
                elif re.match(r"^（[一二三四五六七八九十]+）", line):
                    add_paragraph(doc, line, l2_font, body_size, first_indent=indent)
                elif re.match(r"^\d+\.\s", line):
                    add_paragraph(doc, line, l3_font, body_size, first_indent=indent)
                elif re.match(r"^（\d+）", line):
                    add_paragraph(doc, line, l4_font, body_size, first_indent=indent)
                # 5. List
                # 5. 列表
                elif line.startswith("- ") or line.startswith("* "):
                    add_paragraph(doc, "· " + line[2:], body_font, body_size, first_indent=indent)
                # 6. Attachment
                # 6. 附件
                elif line.startswith("附件") or line.startswith("**附件**"):
                    add_paragraph(doc, line, body_font, body_size, first_indent=indent, space_before=Pt(12))
                # 7. Body text
                # 7. 正文
                else:
                    add_rich_paragraph(doc, line, body_font, body_size, first_indent=indent, bold_font=l1_font)
                i += 1

            # === Page margins ===
            # === 页面边距 ===
            if not is_append:
                for section in doc.sections:
                    section.top_margin = Cm(margins.get("top", 2.54))
                    section.bottom_margin = Cm(margins.get("bottom", 2.54))
                    section.left_margin = Cm(margins.get("left", 3.18))
                    section.right_margin = Cm(margins.get("right", 3.18))

            doc.save(path)
            return f"Successfully created Word document (style: {style_str}). [RESULT_PATH: {os.path.abspath(path)}]"
        except Exception as e:
            return f"Word Write Error: {str(e)}"

    def _resolve_style(self, style_str: str) -> dict:
        """Resolve style parameter: preset name or custom JSON.
        解析样式参数：预设名或自定义 JSON"""
        import json as _json

        # Try as preset name
        # 尝试作为预设名
        if style_str in self.STYLE_PRESETS:
            return dict(self.STYLE_PRESETS[style_str])
        # Try parsing as JSON
        # 尝试作为 JSON 解析
        try:
            custom = _json.loads(style_str)
            if isinstance(custom, dict):
                # Fill in missing fields based on default
                # 基于 default 补全缺失字段
                base = dict(self.STYLE_PRESETS["default"])
                base.update(custom)
                return base
        except (ValueError, TypeError):
            pass
        # Fall back to default
        # 回退到 default
        return dict(self.STYLE_PRESETS["default"])


# --- 3. PDF 工具 ---
class PdfWriteArgs(BaseModel):
    path: str = Field(description="保存 PDF (.pdf) 的路径")
    content: str = Field(description="PDF 的文本内容")
    title: str = Field(description="文档标题", default="Rooster Investigation Report")


class PdfWriteTool(BaseTool):
    """职业化 PDF 生成工具 (基于 fpdf2)"""

    name: str = "office_pdf_write"
    kit: str = "Office"
    fc_hidden: bool = True  # [Round 10] Use pdf_op(action="write") instead
    description: str = "Export text content as a formatted PDF report. Use this to generate deliverable documents."
    domain: str = "craft"
    args_schema: Type[BaseModel] = PdfWriteArgs

    async def run(self, **kwargs) -> str:
        try:
            from fpdf import FPDF
        except ImportError:
            return "Error: 'fpdf2' not installed. Please run 'pip install fpdf2'."

        path = kwargs.get("path")
        content = kwargs.get("content", "")
        title = kwargs.get("title", "Rooster Report")

        try:
            pdf = FPDF()
            pdf.add_page()
            # Try to load a Chinese-compatible font (common Windows paths)
            # 尝试加载能显示中文的字体（Windows 常用路径）
            font_added = False
            for font_p in [
                "C:\\Windows\\Fonts\\simhei.ttf",
                "C:\\Windows\\Fonts\\msyh.ttc",
                "C:\\Windows\\Fonts\\simsun.ttc",
                "/System/Library/Fonts/PingFang.ttc",
                "/System/Library/Fonts/STHeiti Light.ttc",
                "/Library/Fonts/Arial Unicode.ttf",
            ]:
                if os.path.exists(font_p):
                    pdf.add_font("Sans", "", font_p)
                    pdf.add_font("Sans", "B", font_p)  # Reuse for Bold
                    pdf.set_font("Sans", size=12)
                    font_added = True
                    break

            if not font_added:
                pdf.set_font("Helvetica", size=12)

            # Title
            if font_added:
                pdf.set_font("Sans", "B", 16)
            else:
                pdf.set_font("Helvetica", "B", 16)
            pdf.cell(0, 10, title, ln=True, align="C")
            pdf.ln(10)

            # Content
            if font_added:
                pdf.set_font("Sans", "", 12)
            else:
                pdf.set_font("Helvetica", "", 12)
            pdf.multi_cell(0, 10, content)

            pdf.output(path)
            return f"Successfully generated Professional PDF report. [RESULT_PATH: {os.path.abspath(path)}]"
        except Exception as e:
            return f"PDF Write Error: {str(e)}"


class PdfReadArgs(BaseModel):
    path: str = Field(description="PDF 文件路径")


class PdfReadTool(BaseTool):
    """机械化 PDF 读取工具"""

    name: str = "office_pdf_read"
    kit: str = "Office"
    fc_hidden: bool = True  # [Round 10] Use pdf_op(action="read") instead
    description: str = "Extract plain text content from a PDF file. Returns all readable text from the document."
    domain: str = "craft"
    args_schema: Type[BaseModel] = PdfReadArgs

    async def run(self, **kwargs) -> str:
        try:
            from pypdf import PdfReader
        except ImportError:
            return "Error: 'pypdf' not installed. Please run 'pip install pypdf'."

        path = kwargs.get("path")
        try:
            reader = PdfReader(path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text if text.strip() else "Empty PDF or text not extractable."
        except Exception as e:
            return f"PDF Read Error: {str(e)}"


# ---------------------------------------------------------------------------
# [Round 10] excel_op — unified Excel macro
# Replaces: excel_read, excel_write
# ---------------------------------------------------------------------------


class ExcelOpArgs(BaseModel):
    action: str = Field(description="'read' to read a spreadsheet, 'write' to save data as Excel")
    path: str = Field(description="Excel file path (.xlsx)")
    data: Optional[List[Dict[str, Any]]] = Field(
        default=None, description="[write] JSON array of objects, e.g. [{'Company': 'OpenAI', 'Value': '80B'}]"
    )
    sheet_name: Optional[str] = Field(default="Sheet1", description="[write] Sheet name (default: Sheet1)")


class ExcelOpTool(BaseTool):
    """[Round 10] Unified Excel macro: read or write spreadsheets."""

    name: str = "excel_op"
    kit: str = "Office"
    description: str = (
        "Unified Excel tool. Use action='read' to read an existing Excel file and return CSV-style content. "
        "Use action='write' to save structured JSON data as an Excel (.xlsx) file."
    )
    domain: str = "craft"
    args_schema: Type[BaseModel] = ExcelOpArgs

    async def run(self, **kwargs) -> str:
        try:
            import pandas as pd
        except ImportError:
            return "Error: 'pandas' or 'openpyxl' not installed. Please run 'pip install pandas openpyxl'."

        action = kwargs.get("action", "").lower()
        path = kwargs.get("path")

        if action == "read":
            try:
                df = pd.read_excel(path)
                csv_preview = df.to_csv(index=False)
                return f"Excel Content Preview (Total {len(df)} rows):\n{csv_preview}"
            except Exception as e:
                return f"Excel Read Error: {str(e)}"

        elif action == "write":
            data = kwargs.get("data", [])
            sheet_name = kwargs.get("sheet_name", "Sheet1")
            try:
                df = pd.DataFrame(data)
                os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
                df.to_excel(path, index=False, sheet_name=sheet_name)
                return f"Successfully saved {len(data)} rows to Excel. [RESULT_PATH: {os.path.abspath(path)}]"
            except Exception as e:
                return f"Excel Write Error: {str(e)}"

        else:
            return f"Error: Unknown action '{action}'. Valid: 'read', 'write'."


# ---------------------------------------------------------------------------
# [Round 10] pdf_op — unified PDF macro
# Replaces: office_pdf_read, office_pdf_write
# ---------------------------------------------------------------------------


class PdfOpArgs(BaseModel):
    action: str = Field(description="'read' to extract text from a PDF, 'write' to generate a PDF report")
    path: str = Field(description="PDF file path")
    content: Optional[str] = Field(default=None, description="[write] Text content for the PDF")
    title: Optional[str] = Field(
        default="Rooster Investigation Report",
        description="[write] Document title (default: Rooster Investigation Report)",
    )


class PdfOpTool(BaseTool):
    """[Round 10] Unified PDF macro: read or write PDF files."""

    name: str = "pdf_op"
    kit: str = "Office"
    description: str = (
        "Unified PDF tool. Use action='read' to extract all readable text from a PDF file. "
        "Use action='write' to generate a formatted PDF report from text content."
    )
    domain: str = "craft"
    args_schema: Type[BaseModel] = PdfOpArgs

    async def run(self, **kwargs) -> str:
        action = kwargs.get("action", "").lower()
        path = kwargs.get("path")

        if action == "read":
            try:
                from pypdf import PdfReader
            except ImportError:
                return "Error: 'pypdf' not installed. Please run 'pip install pypdf'."
            try:
                reader = PdfReader(path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text if text.strip() else "Empty PDF or text not extractable."
            except Exception as e:
                return f"PDF Read Error: {str(e)}"

        elif action == "write":
            try:
                from fpdf import FPDF
            except ImportError:
                return "Error: 'fpdf2' not installed. Please run 'pip install fpdf2'."
            content = kwargs.get("content", "")
            title = kwargs.get("title", "Rooster Report")
            try:
                pdf = FPDF()
                pdf.add_page()
                font_added = False
                for font_p in [
                    "C:\\Windows\\Fonts\\simhei.ttf",
                    "C:\\Windows\\Fonts\\msyh.ttc",
                    "C:\\Windows\\Fonts\\simsun.ttc",
                    "/System/Library/Fonts/PingFang.ttc",
                    "/System/Library/Fonts/STHeiti Light.ttc",
                    "/Library/Fonts/Arial Unicode.ttf",
                ]:
                    if os.path.exists(font_p):
                        pdf.add_font("Sans", "", font_p)
                        pdf.add_font("Sans", "B", font_p)
                        pdf.set_font("Sans", size=12)
                        font_added = True
                        break
                if not font_added:
                    pdf.set_font("Helvetica", size=12)
                pdf.set_font("Sans" if font_added else "Helvetica", "B", 16)
                pdf.cell(0, 10, title, ln=True, align="C")
                pdf.ln(10)
                pdf.set_font("Sans" if font_added else "Helvetica", "", 12)
                pdf.multi_cell(0, 10, content)
                pdf.output(path)
                return f"Successfully generated PDF report. [RESULT_PATH: {os.path.abspath(path)}]"
            except Exception as e:
                return f"PDF Write Error: {str(e)}"

        else:
            return f"Error: Unknown action '{action}'. Valid: 'read', 'write'."


# ---------------------------------------------------------------------------
# [New] pptx_op — PowerPoint generation tool
# Accepts Markdown-style slide content, handles themes and Chinese fonts.
# ---------------------------------------------------------------------------


class PptxOpArgs(BaseModel):
    path: str = Field(description="输出 .pptx 文件的绝对路径，例如 C:/Users/xxx/Desktop/report.pptx")
    slides_markdown: str = Field(
        description=(
            "幻灯片内容，Markdown 格式，每张幻灯片之间用 '---' 单独一行分隔。\n"
            "每张幻灯片第一行 '# 标题' 作为幻灯片大标题，\n"
            "'## 副标题' 作为副标题，\n"
            "'-' 开头的行作为要点列表。\n"
            "示例：\n"
            "# 第一章 概述\n## 项目背景\n- 要点1\n- 要点2\n---\n# 第二章 方案\n- 核心方案\n- 执行路径"
        )
    )
    theme: str = Field(
        default="business", description="主题风格：'business'（深蓝商务）、'minimal'（白色简约）、'dark'（深色科技）"
    )
    title: str = Field(default="", description="演示文稿全局标题（可选，用于首页封面）")
    author: str = Field(default="", description="作者署名（可选）")
    template_path: Optional[str] = Field(
        default=None, description="自定义 PPTX 模板文件的绝对路径。如果提供，将以此文件为底板生成内容。"
    )


class PptxOpTool(BaseTool):
    """PowerPoint 演示文稿生成工具 — 支持 Markdown 输入和多种主题风格"""

    name: str = "pptx_op"
    kit: str = "Office"
    description: str = (
        "Generate a PowerPoint (.pptx) presentation from Markdown-formatted slide content. "
        "Each slide is separated by '---'. First line '# Title' = slide title, "
        "'## Subtitle' = subtitle, '- item' = bullet point. "
        "Supports themes: 'business' (deep blue), 'minimal' (clean white), 'dark' (tech dark). "
        "Always saves to local disk. Use this instead of writing python-pptx code manually."
    )
    domain: str = "craft"
    args_schema: Type[BaseModel] = PptxOpArgs

    # --- Theme definitions ---
    THEMES = {
        "business": {
            "bg": (0x1F, 0x3D, 0x7A),  # 深海蓝背景
            "title_color": (0xFF, 0xFF, 0xFF),  # 白色标题
            "body_color": (0xCC, 0xDD, 0xFF),  # 浅蓝正文
            "accent": (0xFF, 0xC0, 0x00),  # 金色强调
            "title_font_size": 36,
            "body_font_size": 20,
        },
        "minimal": {
            "bg": (0xFF, 0xFF, 0xFF),  # 白色背景
            "title_color": (0x1A, 0x1A, 0x2E),  # 深蓝标题
            "body_color": (0x33, 0x33, 0x33),  # 深灰正文
            "accent": (0x00, 0x78, 0xD4),  # 蓝色强调
            "title_font_size": 38,
            "body_font_size": 20,
        },
        "dark": {
            "bg": (0x0D, 0x0D, 0x1A),  # 极深蓝背景
            "title_color": (0x00, 0xFF, 0xCC),  # 青绿标题
            "body_color": (0xE0, 0xE0, 0xE0),  # 浅灰正文
            "accent": (0xFF, 0x6B, 0x35),  # 橙色强调
            "title_font_size": 36,
            "body_font_size": 20,
        },
    }

    async def run(self, **kwargs) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return "Error: 'python-pptx' not installed. Please run 'pip install python-pptx'."

        import re

        path = kwargs.get("path", "")
        slides_md = kwargs.get("slides_markdown", "")
        theme_name = kwargs.get("theme", "business")
        prs_title = kwargs.get("title", "")
        author = kwargs.get("author", "")
        template_path = kwargs.get("template_path", "")

        if not path:
            return "Error: 'path' is required."
        if not slides_md.strip():
            return "Error: 'slides_markdown' is required."

        theme = self.THEMES.get(theme_name, self.THEMES["business"])

        try:
            if template_path and os.path.exists(template_path):
                prs = Presentation(template_path)
                use_template = True
                try:
                    cover_layout = prs.slide_layouts[0]
                    content_layout = prs.slide_layouts[1]
                except Exception:
                    cover_layout = prs.slide_layouts[6]
                    content_layout = prs.slide_layouts[6]
            else:
                prs = Presentation()
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                use_template = False
                cover_layout = prs.slide_layouts[6]
                content_layout = prs.slide_layouts[6]

            def rgb(color_tuple):
                return RGBColor(*color_tuple)

            def add_textbox(
                slide,
                text,
                left,
                top,
                width,
                height,
                font_size=20,
                bold=False,
                color=(255, 255, 255),
                align=PP_ALIGN.LEFT,
                word_wrap=True,
            ):
                txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                tf = txBox.text_frame
                tf.word_wrap = word_wrap
                p = tf.paragraphs[0]
                p.alignment = align
                run = p.add_run()
                run.text = text
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = rgb(color)
                # 中文字体设置
                run.font.name = "微软雅黑"
                return txBox

            def set_slide_bg(slide, color_tuple):
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = rgb(color_tuple)

            def add_bullet_slide(slide, title_text, subtitle_text, bullets):
                """渲染一张内容幻灯片"""
                if use_template and slide.shapes.title:
                    # 优先使用模板自带占位符
                    slide.shapes.title.text = title_text

                    body_shape = None
                    for shape in slide.placeholders:
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format.idx == 1:
                            body_shape = shape
                            break

                    # 智能回退：寻找除标题外的任意文本占位符
                    if not body_shape:
                        for shape in slide.placeholders:
                            if shape != slide.shapes.title and hasattr(shape, "text_frame"):
                                body_shape = shape
                                break

                    if body_shape:
                        tf = body_shape.text_frame
                        tf.clear()  # 清空占位符里的提示文字
                        first_p = tf.paragraphs[0]
                        if subtitle_text:
                            first_p.text = subtitle_text
                            # bullets start from new para
                            for bullet in bullets:
                                p = tf.add_paragraph()
                                p.text = bullet
                                p.level = 0
                        else:
                            if bullets:
                                first_p.text = bullets[0]
                                first_p.level = 0
                                for bullet in bullets[1:]:
                                    p = tf.add_paragraph()
                                    p.text = bullet
                                    p.level = 0
                        return

                    # 降级：如果没有占位符，清空已写的 title 避免重叠
                    slide.shapes.title.text = ""

                # 降级：手绘元素
                if not use_template:
                    set_slide_bg(slide, theme["bg"])
                    # 顶部色条
                    accent_bar = slide.shapes.add_shape(
                        1,  # MSO_SHAPE_TYPE.RECTANGLE
                        Inches(0),
                        Inches(0),
                        prs.slide_width,
                        Inches(0.08),
                    )
                    accent_bar.fill.solid()
                    accent_bar.fill.fore_color.rgb = rgb(theme["accent"])
                    accent_bar.line.fill.background()

                # 标题
                add_textbox(
                    slide,
                    title_text,
                    left=0.4,
                    top=0.3,
                    width=12.5,
                    height=1.2,
                    font_size=theme["title_font_size"],
                    bold=True,
                    color=theme["title_color"] if not use_template else (0, 0, 0),
                    align=PP_ALIGN.LEFT,
                )

                # 副标题
                if subtitle_text:
                    add_textbox(
                        slide,
                        subtitle_text,
                        left=0.4,
                        top=1.4,
                        width=12.5,
                        height=0.6,
                        font_size=theme["body_font_size"] - 2,
                        bold=False,
                        color=theme["accent"] if not use_template else (50, 50, 50),
                        align=PP_ALIGN.LEFT,
                    )

                # 要点列表
                if bullets:
                    top_start = 2.1 if subtitle_text else 1.8
                    bullet_height = min(4.8, len(bullets) * 0.65 + 0.3)
                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(top_start), Inches(12.3), Inches(bullet_height)
                    )
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    first = True
                    for bullet in bullets:
                        if first:
                            p = tf.paragraphs[0]
                            first = False
                        else:
                            p = tf.add_paragraph()
                        p.alignment = PP_ALIGN.LEFT
                        p.space_before = Pt(4)
                        run = p.add_run()
                        run.text = f"◆  {bullet}"
                        run.font.size = Pt(theme["body_font_size"])
                        if not use_template:
                            run.font.color.rgb = rgb(theme["body_color"])
                        run.font.name = "微软雅黑"

            # --- 封面幻灯片（如果指定了全局 title）---
            if prs_title:
                cover_slide = prs.slides.add_slide(cover_layout)
                if use_template and cover_slide.shapes.title:
                    cover_slide.shapes.title.text = prs_title
                    author_shape = None
                    for shape in cover_slide.placeholders:
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format.idx == 1:
                            author_shape = shape
                            break
                    if not author_shape:
                        for shape in cover_slide.placeholders:
                            if shape != cover_slide.shapes.title and hasattr(shape, "text_frame"):
                                author_shape = shape
                                break
                    if author_shape:
                        author_shape.text_frame.text = author
                else:
                    if not use_template:
                        set_slide_bg(cover_slide, theme["bg"])
                        # 装饰色块
                        deco = cover_slide.shapes.add_shape(1, Inches(0), Inches(2.8), prs.slide_width, Inches(0.12))
                        deco.fill.solid()
                        deco.fill.fore_color.rgb = rgb(theme["accent"])
                        deco.line.fill.background()

                    # 主标题
                    add_textbox(
                        cover_slide,
                        prs_title,
                        left=0.8,
                        top=1.5,
                        width=11.5,
                        height=1.8,
                        font_size=44,
                        bold=True,
                        color=theme["title_color"] if not use_template else (0, 0, 0),
                        align=PP_ALIGN.CENTER,
                    )
                    if author:
                        add_textbox(
                            cover_slide,
                            author,
                            left=0.8,
                            top=3.2,
                            width=11.5,
                            height=0.6,
                            font_size=18,
                            bold=False,
                            color=theme["body_color"] if not use_template else (50, 50, 50),
                            align=PP_ALIGN.CENTER,
                        )

            # --- 解析 Markdown 幻灯片 ---
            raw_slides = re.split(r"\n---\n|^---$", slides_md, flags=re.MULTILINE)
            slide_count = 0

            for raw in raw_slides:
                raw = raw.strip()
                if not raw:
                    continue

                lines = raw.split("\n")
                title_text = ""
                subtitle_text = ""
                bullets = []
                body_lines = []

                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    if line.startswith("# "):
                        title_text = line[2:].strip()
                    elif line.startswith("## "):
                        subtitle_text = line[3:].strip()
                    elif line.startswith("- ") or line.startswith("* "):
                        bullets.append(line[2:].strip())
                    elif line.startswith("### "):
                        bullets.append(f"【{line[4:].strip()}】")
                    else:
                        body_lines.append(line)

                # 将正文行也转为要点
                for bl in body_lines:
                    if bl and not bl.startswith("|"):  # 跳过表格
                        bullets.append(bl)

                if not title_text and not bullets:
                    continue

                slide = prs.slides.add_slide(content_layout)
                add_bullet_slide(slide, title_text, subtitle_text, bullets)
                slide_count += 1

            if slide_count == 0:
                return "Error: No valid slides found in slides_markdown. Check format: use '# Title' for slide titles and '---' to separate slides."

            os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
            prs.save(path)
            abs_path = os.path.abspath(path)
            return (
                f"Successfully created PowerPoint with {slide_count} slides "
                f"(theme: {theme_name}). [RESULT_PATH: {abs_path}]"
            )

        except Exception as e:
            return f"PPTX Error: {str(e)}"
